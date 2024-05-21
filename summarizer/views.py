from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib.auth.forms import UserCreationForm
from .models import Document, Summary, TextEntry, UrlEntry
from .forms import DocumentForm, SummaryForm, TextForm, UrlForm
import openai
from django.conf import settings
import PyPDF2
from pptx import Presentation
import pandas as pd
import docx  # Add this import for handling Word documents
from django.http import JsonResponse, HttpResponse
import requests
from bs4 import BeautifulSoup
from django.db import models

# Use the OpenAI API key from settings
openai.api_key = settings.OPENAI_API_KEY

def home(request):
    return render(request, 'home.html')

def signup(request):
    if request.method == 'POST':
        form = UserCreationForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('login')
    else:
        form = UserCreationForm()
    return render(request, 'registration/signup.html', {'form': form})

@login_required
def upload_document(request):
    if request.method == 'POST':
        form = DocumentForm(request.POST, request.FILES)
        if form.is_valid():
            document = form.save(commit=False)
            document.user = request.user
            document.save()
            return redirect('summarize_document', document_id=document.id)
    else:
        form = DocumentForm()
    return render(request, 'upload.html', {'form': form})

@login_required
def enter_text(request):
    if request.method == 'POST':
        form = TextForm(request.POST)
        if form.is_valid():
            text_entry = form.save(commit=False)
            text_entry.user = request.user
            text_entry.save()
            return redirect('summarize_text', text_id=text_entry.id)
    else:
        form = TextForm()
    return render(request, 'enter_text.html', {'form': form})

@login_required
def enter_url(request):
    if request.method == 'POST':
        form = UrlForm(request.POST)
        if form.is_valid():
            url_entry = form.save(commit=False)
            url_entry.user = request.user
            url_entry.save()
            return redirect('summarize_url', url_id=url_entry.id)
    else:
        form = UrlForm()
    return render(request, 'enter_url.html', {'form': form})

@login_required
def summarize_document(request, document_id):
    document = get_object_or_404(Document, id=document_id)
    if request.method == 'POST':
        form = SummaryForm(request.POST)
        if form.is_valid():
            word_count = form.cleaned_data['word_count']
            tone = form.cleaned_data['tone'].capitalize()
            summary_text = generate_summary(document.file.path, word_count, tone)
            summary = Summary(document=document, summary_text=summary_text, word_count=word_count, tone=tone)
            summary.save()
            return JsonResponse({'success': True, 'redirect': summary.get_absolute_url()})
        else:
            return JsonResponse({'success': False})
    else:
        form = SummaryForm()
    return render(request, 'summarize.html', {'form': form, 'document': document})

@login_required
def summarize_text(request, text_id):
    text_entry = get_object_or_404(TextEntry, id=text_id)
    if request.method == 'POST':
        form = SummaryForm(request.POST)
        if form.is_valid():
            word_count = form.cleaned_data['word_count']
            tone = form.cleaned_data['tone'].capitalize()
            summary_text = generate_text_summary(text_entry.text, word_count, tone)
            summary = Summary(text_entry=text_entry, summary_text=summary_text, word_count=word_count, tone=tone)
            summary.save()
            return JsonResponse({'success': True, 'redirect': summary.get_absolute_url()})
        else:
            return JsonResponse({'success': False})
    else:
        form = SummaryForm()
    return render(request, 'summarize.html', {'form': form, 'text_entry': text_entry})

@login_required
def summarize_url(request, url_id):
    url_entry = get_object_or_404(UrlEntry, id=url_id)
    if request.method == 'POST':
        form = SummaryForm(request.POST)
        if form.is_valid():
            word_count = form.cleaned_data['word_count']
            tone = form.cleaned_data['tone'].capitalize()
            article_text = fetch_url_text(url_entry.url)
            summary_text = generate_text_summary(article_text, word_count, tone)
            summary = Summary(url_entry=url_entry, summary_text=summary_text, word_count=word_count, tone=tone)
            summary.save()
            return JsonResponse({'success': True, 'redirect': summary.get_absolute_url()})
        else:
            return JsonResponse({'success': False})
    else:
        form = SummaryForm()
    return render(request, 'summarize.html', {'form': form, 'url_entry': url_entry})

def extract_text_from_pdf(file_path):
    try:
        text = ""
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text()
        return text
    except PdfReadError:
        return "The uploaded PDF file is corrupted or cannot be read."

def extract_text_from_pptx(file_path):
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_text_from_excel(file_path):
    text = ""
    xls = pd.ExcelFile(file_path)
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        text += df.to_string(index=False) + " "
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + " "
    return text

def extract_text_from_txt(file_path):
    with open(file_path, 'r') as file:
        text = file.read()
    return text

def generate_summary(file_path, word_count, tone):
    if file_path.lower().endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
    elif file_path.lower().endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    elif file_path.lower().endswith('.xlsx'):
        text = extract_text_from_excel(file_path)
    elif file_path.lower().endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.lower().endswith('.txt'):
        text = extract_text_from_txt(file_path)
    else:
        return "Unsupported file type."
    
    if text.startswith("The uploaded PDF file is corrupted") or text == "Unsupported file type.":
        return text
    
    # Generate summary using OpenAI API
    response = openai.ChatCompletion.create(
        model="gpt-4-turbo",
        messages=[
            {"role": "system", "content": f"Summarize the following text in {word_count} words with a {tone} tone."},
            {"role": "user", "content": text}
        ],
        max_tokens=1500
    )
    summary = response.choices[0].message['content'].strip()
    return summary

def generate_text_summary(text, word_count, tone):
    # Generate summary using OpenAI API
    response = openai.ChatCompletion.create(
        model="gpt-4-turbo",
        messages=[
            {"role": "system", "content": f"Summarize the following text in {word_count} words with a {tone} tone."},
            {"role": "user", "content": text}
        ],
        max_tokens=1500
    )
    summary = response.choices[0].message['content'].strip()
    return summary

def fetch_url_text(url):
    # Implement the logic to fetch and extract text from the URL
    response = requests.get(url)
    soup = BeautifulSoup(response.content, 'html.parser')
    paragraphs = soup.find_all('p')
    article_text = ' '.join([paragraph.get_text() for paragraph in paragraphs])
    return article_text

@login_required
def history(request):
    summaries = Summary.objects.filter(
        models.Q(document__user=request.user) |
        models.Q(text_entry__user=request.user) |
        models.Q(url_entry__user=request.user)
    ).order_by('-created_at')
    return render(request, 'history.html', {'summaries': summaries})

@login_required
def summary_detail(request, summary_id):
    summary = get_object_or_404(Summary, id=summary_id)
    return render(request, 'summary_detail.html', {'summary': summary})

@login_required
def delete_summary(request, summary_id):
    summary = get_object_or_404(Summary, id=summary_id)
    if summary.document and summary.document.user == request.user or \
       summary.text_entry and summary.text_entry.user == request.user or \
       summary.url_entry and summary.url_entry.user == request.user:
        summary.delete()
        return HttpResponse(status=204)
    return HttpResponse(status=403)

@login_required
def delete_history(request):
    Summary.objects.filter(
        models.Q(document__user=request.user) |
        models.Q(text_entry__user=request.user) |
        models.Q(url_entry__user=request.user)
    ).delete()
    return HttpResponse(status=204)
